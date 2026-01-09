"""Generate a 10-slide pitch deck using python-pptx.

Usage:
  pip install python-pptx
  python scripts/generate_pitch_deck.py

Output: docs/pitch_deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

slides = [
    ("InsiderGuard — Predict Insider Threats Before Data Loss", "One-line elevator pitch: Behavioral baselines + cognitive intent + HMM phase inference to surface high-confidence insider-threat alerts before exfiltration."),
    ("Problem & Opportunity", "- Data egress and insider risk are hard to detect\n- Existing rules generate noise; security teams overwhelmed\n- Opportunity: early, high-precision signals to prevent data loss"),
    ("Solution Overview", "- Train per-user, peer, and temporal baselines\n- Real-time anomaly scoring + intent engine + HMM phases\n- Dashboard & API for investigation and response"),
    ("Architecture (High Level)", "Data Ingestion → Feature Engineering → Distribution Builder → Scoring → Alerts\nStage-3: Cognitive intent (parso) → Stage-4: HMM phase inference"),
    ("Model & Algorithms", "- Distribution-based personal/peer/temporal models (gaussian/percentile)\n- Fusion weighting (personal, peer, temporal)\n- Stage-3: EMA cognitive updates\n- Stage-4: Gaussian HMM online inference"),
    ("Demo & Results", "- Integration script: train + Stage-3 & Stage-4 demo modes\n- Local demo shows stage progression and high-confidence alerts\n- Tests and smoke integration present"),
    ("Integration & Deployment", "- Flask API + React frontend\n- Docker + docker-compose for local and dev runs\n- Bicep infra templates for App Insights (Azure)"),
    ("Security & Privacy", "- PII hashed and salted, optional E2E for sensitive sources\n- Configurable privacy filters and encryption badges in onboarding"),
    ("Business Metrics & ROI", "- Reduction in incident dwell time\n- Fewer false positives via behavioral baselines\n- Faster triage with confidence scores and explanations"),
    ("Next Steps & Ask", "- Pilot with 1–2 teams\n- Add Kafka integration for Stage-3 production mode\n- CI + production deployment and monitoring (App Insights)")
]

out = Path('docs')
out.mkdir(parents=True, exist_ok=True)
prs = Presentation()
for title, body in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tx = slide.shapes.placeholders[1].text_frame
    for i, line in enumerate(body.split('\n')):
        p = tx.add_paragraph() if i > 0 else tx.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(14)

prs.save(out / 'pitch_deck.pptx')
print('Generated docs/pitch_deck.pptx')
